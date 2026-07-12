"""V6-10/V6-11: PPTX layout parsing tests.

Tests use the real ``operating-system-slides.pptx`` fixture (6 slides,
multi-level bullets, a table, shapes, and an image-only slide).
"""
import os

from app.retrieval.parsers import parse_pptx
from app.retrieval.document_ir import (
    DocumentBlock,
    DocumentPage,
    to_document_pages,
)

FIXTURES = os.path.join(os.path.dirname(__file__), "fixtures", "course_materials")
PPTX_PATH = os.path.join(FIXTURES, "operating-system-slides.pptx")


def _parse_to_doc():
    """Parse the fixture PPTX and convert every slide to the V6 document IR."""
    pages = parse_pptx(PPTX_PATH)
    return to_document_pages(pages)


# ---------------------------------------------------------------------------
# V6-11: multi-level bullet hierarchy preserved
# ---------------------------------------------------------------------------

def test_multi_level_bullet_hierarchy_preserved():
    """Paragraph ``level`` must be carried into ``list_level`` on blocks.

    Slide 2 has level-0 paragraphs ("进程调度概述", "进程调度基础",
    "TLB 与地址转换") and level-1 paragraphs ("FCFS: 先来先服务", ...).
    """
    doc = _parse_to_doc()
    slide2 = doc[1]  # 0-indexed

    # At least one block should have list_level >= 1.
    level1_blocks = [b for b in slide2.blocks if b.list_level >= 1]
    assert level1_blocks, (
        "expected at least one level-1 bullet on slide 2, got levels: "
        f"{[b.list_level for b in slide2.blocks]}"
    )

    # "FCFS" is a level-1 bullet.
    fcfs = [b for b in slide2.blocks if "FCFS" in b.text]
    assert fcfs, "expected 'FCFS' block on slide 2"
    assert fcfs[0].list_level >= 1, f"FCFS should be level 1, got {fcfs[0].list_level}"


# ---------------------------------------------------------------------------
# V6-11: shape sorting by position
# ---------------------------------------------------------------------------

def test_shape_sorting_by_position():
    """Shapes must be sorted by (top, left) for correct reading order.

    Slide 4 has three auto-shapes at different positions:
      - "中断驱动 I/O"    at top=1828800, left=914400
      - "DMA 传输"        at top=1828800, left=4572000
      - "通道控制方式"    at top=3657600, left=914400
    Reading order: 中断驱动 I/O -> DMA 传输 -> 通道控制方式.
    """
    doc = _parse_to_doc()
    slide4 = doc[3]  # 0-indexed

    texts = [b.text for b in slide4.blocks if b.block_type == "body"]
    assert "中断驱动 I/O" in texts
    assert "DMA 传输" in texts
    assert "通道控制方式" in texts

    # 中断驱动 I/O (top=1828800, left=914400) before DMA (top=1828800, left=4572000)
    idx_zd = texts.index("中断驱动 I/O")
    idx_dma = texts.index("DMA 传输")
    idx_td = texts.index("通道控制方式")
    assert idx_zd < idx_dma < idx_td, (
        f"shape order wrong: {texts}"
    )


# ---------------------------------------------------------------------------
# V6-11: table content extraction
# ---------------------------------------------------------------------------

def test_table_content_extraction():
    """Table cells must be extracted as ``table`` blocks with all content.

    Slide 3 has a 3-column x 4-row table:
      header: 策略 | 优点 | 缺点
      row 1:  分页 | 无外部碎片 | 有内部碎片
      row 2:  分段 | 逻辑清晰   | 有外部碎片
      row 3:  段页式 | 综合优点  | 地址转换复杂
    """
    doc = _parse_to_doc()
    slide3 = doc[2]  # 0-indexed

    table_blocks = [b for b in slide3.blocks if b.block_type == "table"]
    table_text = " ".join(b.text for b in table_blocks)

    assert "策略" in table_text, f"expected table header '策略', got {table_text!r}"
    assert "分页" in table_text, f"expected '分页' row, got {table_text!r}"
    assert "分段" in table_text, f"expected '分段' row, got {table_text!r}"
    assert "段页式" in table_text, f"expected '段页式' row, got {table_text!r}"


# ---------------------------------------------------------------------------
# V6-11: image-only slide
# ---------------------------------------------------------------------------

def test_image_only_slide_marked():
    """Slide 5 (only a figure caption, no title) must be ``image_only``."""
    doc = _parse_to_doc()
    slide5 = doc[4]  # 0-indexed

    assert slide5.page_type == "image_only", (
        f"slide 5 should be image_only, got {slide5.page_type!r}"
    )


# ---------------------------------------------------------------------------
# V6-11: single-letter shape labels not treated as body
# ---------------------------------------------------------------------------

def test_single_letter_shape_label_not_body(tmp_path):
    """A shape whose text is a single letter must not become a body block."""
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    from pptx.util import Inches

    txbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    txbox.text_frame.text = "A"  # single-letter label

    path = str(tmp_path / "single.pptx")
    prs.save(path)

    pages = parse_pptx(path)
    doc = to_document_pages(pages)
    body_texts = [b.text for b in doc[0].blocks if b.block_type == "body"]
    assert "A" not in body_texts, (
        "single-letter shape label should not be treated as body"
    )


# ---------------------------------------------------------------------------
# V6-10: title placeholders marked as heading
# ---------------------------------------------------------------------------

def test_title_placeholder_marked_as_heading():
    """Title placeholder text must be classified as ``heading``."""
    doc = _parse_to_doc()
    slide2 = doc[1]

    title_blocks = [b for b in slide2.blocks if b.block_type == "heading"]
    assert any("进程调度算法" in b.text for b in title_blocks), (
        f"expected title '进程调度算法' as heading, got {[b.text for b in title_blocks]}"
    )


# ---------------------------------------------------------------------------
# V6-10: document IR for PPTX
# ---------------------------------------------------------------------------

def test_pptx_document_page_structure():
    """Each slide must produce a DocumentPage with layout-v6 version."""
    doc = _parse_to_doc()
    assert len(doc) == 6, f"expected 6 slides, got {len(doc)}"
    for page in doc:
        assert isinstance(page, DocumentPage)
        assert page.parser_version == "layout-v6"
        for block in page.blocks:
            assert isinstance(block, DocumentBlock)
            assert block.source_kind == "pptx"
