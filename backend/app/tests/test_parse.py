"""Tests for the material parsing and chunking module.

Strict TDD: these tests are written first and fail until the parse
endpoint, chunker, parsers, and MaterialChunk model are implemented.

Covers:
- POST /api/v1/materials/{id}/parse (txt, pdf, corrupt)
- GET  /api/v1/materials/{id}/chunks (paginated)
- 404 on non-existent material
- 404 on cross-user access
- Status transitions (uploaded -> processing -> ready/failed)
- Unit tests for chunk_text strategy and parsers
"""
import io

import pytest

from app.core.exceptions import BusinessException
from app.retrieval.chunker import (
    build_chunks,
    chunk_text,
    clean_keyword_text,
)
from app.retrieval.parsers import parse_file, parse_pdf, parse_txt
from app.tests.conftest import auth_headers, create_course, upload_material


# A long Chinese paragraph (~3000 chars) used to exercise the chunker.
LONG_TEXT = (
    "操作系统是计算机系统中最基本的系统软件，"
    "它负责管理计算机的硬件资源和软件资源，"
    "为用户和应用程序提供一个方便的接口。"
    "操作系统的基本功能包括进程管理、内存管理、文件系统、设备管理和用户接口。"
    "进程管理负责创建、调度和销毁进程，保证多个进程公平地使用 CPU。"
    "内存管理负责分配和回收内存空间，提供虚拟内存机制。"
    "文件系统负责存储和管理文件，提供目录结构和访问控制。"
    "设备管理负责管理输入输出设备，为上层提供统一的接口。\n"
) * 30


def test_parse_txt_material(client, tmp_path, monkeypatch) -> None:
    """POST /materials/{id}/parse on a .txt returns ready + chunks."""
    monkeypatch.setattr("app.core.config.settings.UPLOAD_DIR", str(tmp_path))
    monkeypatch.setattr(
        "app.core.config.settings.PARSED_DIR", str(tmp_path / "parsed")
    )

    headers = auth_headers(client, username="alice")
    course_id = create_course(client, headers, "操作系统")
    material_id = upload_material(
        client,
        headers,
        course_id,
        "notes.txt",
        LONG_TEXT.encode("utf-8"),
    )

    parse_resp = client.post(
        f"/api/v1/materials/{material_id}/parse", headers=headers
    )
    assert parse_resp.status_code == 200
    body = parse_resp.json()
    assert body["material_id"] == material_id
    assert body["status"] == "ready"
    assert body["chunk_count"] > 0

    chunks_resp = client.get(
        f"/api/v1/materials/{material_id}/chunks", headers=headers
    )
    assert chunks_resp.status_code == 200
    chunks_body = chunks_resp.json()
    items = chunks_body["items"] if isinstance(chunks_body, dict) else chunks_body
    assert len(items) > 0
    for chunk in items:
        assert "text" in chunk
        assert "chunk_index" in chunk


def test_parse_pdf_material(client, tmp_path, monkeypatch) -> None:
    """POST /materials/{id}/parse on a pypdf-generated PDF returns ready."""
    try:
        from pypdf import PdfWriter
    except ImportError:
        pytest.skip("pypdf not available")

    monkeypatch.setattr("app.core.config.settings.UPLOAD_DIR", str(tmp_path))
    monkeypatch.setattr(
        "app.core.config.settings.PARSED_DIR", str(tmp_path / "parsed")
    )

    # Generate a minimal valid PDF (blank page) so the parser succeeds.
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    pdf_buffer = io.BytesIO()
    writer.write(pdf_buffer)
    pdf_bytes = pdf_buffer.getvalue()

    headers = auth_headers(client, username="alice")
    course_id = create_course(client, headers, "操作系统")
    material_id = upload_material(
        client, headers, course_id, "doc.pdf", pdf_bytes
    )

    parse_resp = client.post(
        f"/api/v1/materials/{material_id}/parse", headers=headers
    )
    assert parse_resp.status_code == 200
    assert parse_resp.json()["status"] == "ready"


def test_parse_invalid_material_id(client) -> None:
    """POST /materials/{id}/parse with unknown id returns 404."""
    headers = auth_headers(client, username="alice")
    response = client.post(
        "/api/v1/materials/99999/parse", headers=headers
    )
    assert response.status_code == 404


def test_parse_other_user_material(client, tmp_path, monkeypatch) -> None:
    """User B parsing user A's material returns 404 (isolation)."""
    monkeypatch.setattr("app.core.config.settings.UPLOAD_DIR", str(tmp_path))
    monkeypatch.setattr(
        "app.core.config.settings.PARSED_DIR", str(tmp_path / "parsed")
    )

    headers_a = auth_headers(client, username="alice")
    course_id = create_course(client, headers_a, "操作系统")
    material_id = upload_material(
        client, headers_a, course_id, "notes.txt", b"hello world"
    )

    headers_b = auth_headers(client, username="bob")
    response = client.post(
        f"/api/v1/materials/{material_id}/parse", headers=headers_b
    )
    assert response.status_code == 404


def test_chunks_pagination(client, tmp_path, monkeypatch) -> None:
    """GET /materials/{id}/chunks supports page/page_size pagination."""
    monkeypatch.setattr("app.core.config.settings.UPLOAD_DIR", str(tmp_path))
    monkeypatch.setattr(
        "app.core.config.settings.PARSED_DIR", str(tmp_path / "parsed")
    )

    headers = auth_headers(client, username="alice")
    course_id = create_course(client, headers, "操作系统")
    material_id = upload_material(
        client,
        headers,
        course_id,
        "notes.txt",
        LONG_TEXT.encode("utf-8"),
    )
    client.post(f"/api/v1/materials/{material_id}/parse", headers=headers)

    resp = client.get(
        f"/api/v1/materials/{material_id}/chunks?page=1&page_size=5",
        headers=headers,
    )
    assert resp.status_code == 200
    body = resp.json()
    assert isinstance(body, dict)
    assert "items" in body
    assert "total" in body
    assert "page" in body
    assert "page_size" in body
    assert body["page"] == 1
    assert body["page_size"] == 5
    assert len(body["items"]) <= 5


def test_parse_failed_status(client, tmp_path, monkeypatch) -> None:
    """Parsing a corrupt file sets status=failed with error_message."""
    monkeypatch.setattr("app.core.config.settings.UPLOAD_DIR", str(tmp_path))
    monkeypatch.setattr(
        "app.core.config.settings.PARSED_DIR", str(tmp_path / "parsed")
    )

    headers = auth_headers(client, username="alice")
    course_id = create_course(client, headers, "操作系统")
    # Upload bytes that look like a PDF by extension but are not a real PDF.
    material_id = upload_material(
        client, headers, course_id, "fake.pdf", b"not a real pdf"
    )

    response = client.post(
        f"/api/v1/materials/{material_id}/parse", headers=headers
    )
    assert response.status_code == 200
    body = response.json()
    assert body["status"] == "failed"
    assert body["chunk_count"] == 0


# ---------------------------------------------------------------------------
# Phase 2 bugfix P1-3: parse rollback safety
# ---------------------------------------------------------------------------


def test_parse_scanner_failure_rolls_back_and_preserves_old_chunks(
    client, tmp_path, monkeypatch
) -> None:
    """When the security scanner raises during re-parse, the transaction
    must roll back so the previously-saved chunks are NOT replaced by
    half-committed new ones.

    Without rollback the except block commits the DELETE of old chunks
    plus the INSERT of new (flushed) chunks — destroying the last known
    good state and leaving the material "failed" with new chunks that
    have no security findings.
    """
    monkeypatch.setattr("app.core.config.settings.UPLOAD_DIR", str(tmp_path))
    monkeypatch.setattr(
        "app.core.config.settings.PARSED_DIR", str(tmp_path / "parsed")
    )

    headers = auth_headers(client, username="alice")
    course_id = create_course(client, headers, "操作系统")
    material_id = upload_material(
        client, headers, course_id, "notes.txt", LONG_TEXT.encode("utf-8")
    )

    # First parse succeeds → material is ready with chunks.
    first = client.post(
        f"/api/v1/materials/{material_id}/parse", headers=headers
    )
    assert first.status_code == 200
    assert first.json()["status"] == "ready"
    first_chunks = client.get(
        f"/api/v1/materials/{material_id}/chunks",
        params={"page": 1, "page_size": 100},
        headers=headers,
    )
    original_texts = [c["text"][:40] for c in first_chunks.json()["items"]]
    assert len(original_texts) > 0

    # Now monkeypatch build_chunks to produce DIFFERENT text on re-parse,
    # so we can tell old chunks from new ones even when SQLite reuses rowids.
    from app.retrieval import chunker as chunker_mod

    real_build = chunker_mod.build_chunks
    call_count = {"n": 0}

    def marker_build(pages, chunk_size=600, overlap=100):
        call_count["n"] += 1
        chunks = real_build(pages, chunk_size, overlap)
        # marker_build is only installed AFTER the first parse, so every
        # call here is a re-parse. Prefix every chunk's text so we can
        # detect whether the new (uncommitted-then-rolled-back) chunks
        # leaked into the DB.
        for c in chunks:
            c["text"] = "MODIFIED_" + c["text"]
        return chunks

    monkeypatch.setattr(
        "app.services.material_parser.build_chunks", marker_build
    )

    # Scanner raises during the re-parse.
    def boom(*args, **kwargs):
        raise RuntimeError("scanner exploded")

    monkeypatch.setattr(
        "app.services.security_scanner.scan_material_chunks", boom
    )

    second = client.post(
        f"/api/v1/materials/{material_id}/parse", headers=headers
    )
    assert second.status_code == 200
    # P0: when old chunks exist, a failed re-parse must keep status="ready"
    # so the previous parse result stays visible to the user.
    assert second.json()["status"] == "ready"

    # After the failed re-parse, the chunks must still be the ORIGINAL
    # ones (no "MODIFIED_" prefix). If rollback is missing, the new
    # "MODIFIED_" chunks would be committed and visible here.
    after = client.get(
        f"/api/v1/materials/{material_id}/chunks",
        params={"page": 1, "page_size": 100},
        headers=headers,
    )
    after_texts = [c["text"][:40] for c in after.json()["items"]]
    assert after_texts == original_texts, (
        "re-parse failure must preserve the original chunks "
        "(rollback should undo the half-finished delete+insert)"
    )
    assert not any(t.startswith("MODIFIED_") for t in after_texts), (
        "new (uncommitted) chunks leaked into the DB — missing rollback"
    )

    # Verify the material row reflects "ready with stale-result warning":
    # status stays ready (old chunks still usable) but error_message
    # records that the latest re-parse failed.
    list_resp = client.get(
        f"/api/v1/courses/{course_id}/materials", headers=headers
    )
    items = list_resp.json()["items"]
    mat_row = next(m for m in items if m["id"] == material_id)
    assert mat_row["status"] == "ready"
    assert mat_row["error_message"]
    assert "上一版" in mat_row["error_message"] or "上次" in mat_row["error_message"]


def test_parse_failure_without_old_chunks_still_failed(
    client, tmp_path, monkeypatch
) -> None:
    """P0: when there are NO old chunks (first parse fails), status must
    stay "failed" — there is no previous result to fall back to.
    """
    monkeypatch.setattr("app.core.config.settings.UPLOAD_DIR", str(tmp_path))
    monkeypatch.setattr(
        "app.core.config.settings.PARSED_DIR", str(tmp_path / "parsed")
    )

    headers = auth_headers(client, username="alice")
    course_id = create_course(client, headers, "操作系统")
    # Corrupt PDF: first parse, no prior chunks exist.
    material_id = upload_material(
        client, headers, course_id, "fake.pdf", b"not a real pdf"
    )

    resp = client.post(
        f"/api/v1/materials/{material_id}/parse", headers=headers
    )
    assert resp.status_code == 200
    body = resp.json()
    assert body["status"] == "failed"
    assert body["chunk_count"] == 0

    list_resp = client.get(
        f"/api/v1/courses/{course_id}/materials", headers=headers
    )
    mat_row = next(
        m for m in list_resp.json()["items"] if m["id"] == material_id
    )
    assert mat_row["status"] == "failed"
    assert mat_row["error_message"]


def test_chunk_size_strategy() -> None:
    """Unit test: chunk_text on long text produces overlapping chunks."""
    text = "操作系统" * 500  # 2000 chars
    chunks = chunk_text(text, chunk_size=600, overlap=100)
    assert len(chunks) > 1
    # Each non-final chunk should be in the 500-800 range (chunk_size=600).
    for chunk in chunks[:-1]:
        assert 500 <= len(chunk["text"]) <= 800
    # Overlap: the second chunk starts with the last 100 chars of the first.
    if len(chunks) >= 2:
        overlap_text = chunks[0]["text"][-100:]
        assert chunks[1]["text"].startswith(overlap_text)
    # Each chunk has the required fields and sequential indices.
    for i, chunk in enumerate(chunks):
        assert chunk["chunk_index"] == i
        assert "text" in chunk
        assert "title" in chunk
        assert "page_no" in chunk


def test_parse_txt_unit(tmp_path) -> None:
    """parse_txt reads a UTF-8 text file."""
    path = tmp_path / "sample.txt"
    path.write_text("Hello 操作系统", encoding="utf-8")
    assert parse_txt(str(path)) == "Hello 操作系统"


def test_parse_pdf_unit(tmp_path) -> None:
    """parse_pdf returns [(page_no, text)] for a pypdf-generated PDF."""
    try:
        from pypdf import PdfWriter
    except ImportError:
        pytest.skip("pypdf not available")
    path = tmp_path / "blank.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    with open(path, "wb") as f:
        writer.write(f)
    pages = parse_pdf(str(path))
    assert isinstance(pages, list)
    assert len(pages) == 1
    page_no, text = pages[0]
    assert isinstance(page_no, int)
    assert isinstance(text, str)


def test_build_chunks_preserves_page_no() -> None:
    """build_chunks propagates page_no from parsed pages."""
    pages = [(1, "甲" * 1500), (2, "乙" * 200)]
    chunks = build_chunks(pages, chunk_size=600, overlap=100)
    assert len(chunks) > 0
    # All chunks carry a page_no from the originating page.
    for chunk in chunks:
        assert chunk["page_no"] in (1, 2)
    # Chunk indices are globally sequential.
    indices = [c["chunk_index"] for c in chunks]
    assert indices == list(range(len(chunks)))


def test_clean_keyword_text_basic() -> None:
    """clean_keyword_text collapses whitespace."""
    raw = "操作系统\n\n  进程  管理\t\n 内存管理"
    cleaned = clean_keyword_text(raw)
    assert "\n" not in cleaned
    assert "\t" not in cleaned
    assert "  " not in cleaned
    assert "操作系统" in cleaned


def test_parse_md_file(tmp_path) -> None:
    """parse_file on a .md file returns [(None, text)] with full content."""
    content = (
        "# 操作系统笔记\n"
        "\n"
        "## 进程管理\n"
        "\n"
        "- 进程调度\n"
        "- 内存分配\n"
        "\n"
        "```python\n"
        "def hello():\n"
        "    print('hi')\n"
        "```\n"
    )
    path = tmp_path / "notes.md"
    path.write_text(content, encoding="utf-8")

    pages = parse_file(str(path), "md")
    assert pages == [(None, content)]


def test_parse_pptx_file(tmp_path) -> None:
    """parse_file on a .pptx returns [(page_no, text)] per slide."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        pytest.skip("python-pptx not available")

    prs = Presentation()
    # Slide 1: title + body text box
    slide1_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(slide1_layout)
    slide1.shapes.title.text = "操作系统"
    slide1.shapes.placeholders[1].text = "进程管理"

    # Slide 2: blank layout with a text box
    slide2_layout = prs.slide_layouts[5]
    slide2 = prs.slides.add_slide(slide2_layout)
    slide2.shapes.title.text = "内存管理"
    textbox = slide2.shapes.add_textbox(
        Inches(1), Inches(1), Inches(4), Inches(2)
    )
    textbox.text_frame.text = "虚拟内存"

    path = tmp_path / "slides.pptx"
    prs.save(str(path))

    pages = parse_file(str(path), "pptx")
    assert len(pages) == 2
    assert pages[0][0] == 1
    assert "操作系统" in pages[0][1]
    assert "进程管理" in pages[0][1]
    assert pages[1][0] == 2
    assert "内存管理" in pages[1][1]
    assert "虚拟内存" in pages[1][1]


def test_parse_file_md_dispatch(tmp_path) -> None:
    """parse_file dispatches by file_type='md' and returns [(None, text)]."""
    path = tmp_path / "notes.md"
    path.write_text("# 标题\n正文", encoding="utf-8")

    pages = parse_file(str(path), "md")
    assert pages == [(None, "# 标题\n正文")]


def test_parse_file_pptx_dispatch(tmp_path) -> None:
    """parse_file dispatches by file_type='pptx' and returns paginated list."""
    try:
        from pptx import Presentation
    except ImportError:
        pytest.skip("python-pptx not available")

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "封面"

    path = tmp_path / "slides.pptx"
    prs.save(str(path))

    pages = parse_file(str(path), "pptx")
    assert isinstance(pages, list)
    assert len(pages) == 1
    page_no, text = pages[0]
    assert page_no == 1
    assert "封面" in text


def test_parse_file_unsupported_still_raises(tmp_path) -> None:
    """parse_file raises BusinessException for unsupported file types."""
    path = tmp_path / "data.xlsx"
    path.write_bytes(b"fake xlsx bytes")

    with pytest.raises(BusinessException):
        parse_file(str(path), "xlsx")
