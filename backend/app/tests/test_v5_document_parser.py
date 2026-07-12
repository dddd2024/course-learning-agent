from app.retrieval.parsers import parse_pdf, parse_pptx
from app.services.material_cleaner import clean_pages


def test_pdf_layout_parser_preserves_page_order_and_footer_cleaning(tmp_path):
    from pypdf import PdfWriter
    from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject
    path = tmp_path / "sample.pdf"; writer = PdfWriter()
    font = writer._add_object(DictionaryObject({NameObject("/Type"): NameObject("/Font"), NameObject("/Subtype"): NameObject("/Type1"), NameObject("/BaseFont"): NameObject("/Helvetica")}))
    for text in ("TCP/IP protocol", "CSMA/CD control"):
        page = writer.add_blank_page(width=612, height=792)
        page[NameObject("/Resources")] = DictionaryObject({NameObject("/Font"): DictionaryObject({NameObject("/F1"): font})})
        stream = DecodedStreamObject(); stream.set_data(f"BT /F1 18 Tf 72 720 Td ({text}) Tj ET".encode())
        page[NameObject("/Contents")] = writer._add_object(stream)
    with path.open("wb") as output: writer.write(output)
    pages = parse_pdf(str(path))
    cleaned = clean_pages([page.text for page in pages])
    assert [page.page_no for page in pages] == [1, 2]
    assert pages[0].parser_version == "pypdf-legacy"
    assert "TCP/IP" in cleaned[0].text
    assert clean_pages(["课程页眉\n正文\n1", "课程页眉\n第二页\n2"])[0].text == "正文"


def test_pptx_parser_uses_shape_order_and_title_placeholder(tmp_path):
    from pptx import Presentation
    path = tmp_path / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "传输层"
    slide.placeholders[1].text = "可靠传输\n流量控制"
    presentation.save(path)
    pages = parse_pptx(str(path))
    assert pages[0].page_no == 1
    assert pages[0].blocks[0].text == "传输层"
    assert pages[0].blocks[0].block_type == "title"


def test_cleaner_marks_graphic_labels_without_dropping_protocol_terms():
    page = clean_pages(["图 1\nTCP/IP\nCSMA/CD\nhttps://bad.example\n正文"])[0]
    assert page.text.splitlines() == ["TCP/IP", "CSMA/CD", "正文"]
    assert {row["reason"] for row in page.decisions if row["decision"] == "removed"} == {"graphic_or_animation_label", "standalone_url"}
