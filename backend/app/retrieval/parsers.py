"""File parsers that extract text content from uploaded materials.

Each parser returns text in a uniform shape so the chunker can consume
it without knowing the source format:

* ``parse_txt`` / ``parse_docx`` return a single string (no page concept)
* ``parse_pdf`` returns ``[(page_no, text), ...]`` (one entry per page)
* ``parse_file`` dispatches by extension and normalises the output to a
  list of ``(page_no, text)`` tuples where ``page_no`` is ``None`` for
  non-paginated formats.

V6-11 improvements (PDF):
  - Two-column layout detection (blocks cluster around two x-ranges).
    Left column is read top-to-bottom first, then right column.
  - Heading detection via font-size statistics (median + std), not a
    fixed threshold.
  - Footer detection: bottom 10% of page AND cross-page repetition.
  - Table detection: blocks within a vector-drawing grid are classified
    as ``table``.
  - Image anchors: filled-shape drawings and raster images are recorded
    with page_no and bbox.

V6-11 improvements (PPTX):
  - Shapes sorted by ``(top, left)`` for reading order.
  - Paragraph ``level`` preserved as ``list_level`` on blocks.
  - Title placeholders, body text, tables, and shapes marked separately.
  - Single-letter shape labels are not treated as body knowledge points.
  - Image-only slides (no title, only captions) are detected.

Errors from the underlying libraries (e.g. ``PdfReadError`` on a corrupt
PDF) propagate to the caller so the parse endpoint can mark the material
as ``failed`` with the exception message.
"""
from __future__ import annotations

import statistics
from dataclasses import dataclass, field
from typing import Any, List, Optional, Tuple

from app.core.exceptions import BusinessException


@dataclass(frozen=True)
class TextBlock:
    text: str
    bbox: tuple[float, float, float, float] = (0, 0, 0, 0)
    font_size: float | None = None
    block_type: str = "body"
    reading_order: int = 0
    source_kind: str = "text"
    list_level: int = 0


@dataclass(frozen=True)
class ParsedPage:
    page_no: int | None
    blocks: list[TextBlock] = field(default_factory=list)
    source_kind: str = "text"
    parser_version: str = "layout-v5"
    # V6: image anchors and page geometry carried for IR conversion.
    images: list = field(default_factory=list)
    page_height: float = 0.0

    @property
    def text(self) -> str:
        return "\n".join(block.text for block in self.blocks if block.text.strip())

    def __iter__(self):
        """Backward-compatible ``(page_no, text)`` unpacking for callers."""
        yield self.page_no
        yield self.text

    def __getitem__(self, index: int):
        return (self.page_no, self.text)[index]

    def __len__(self) -> int:
        return 2


# ---------------------------------------------------------------------------
# Plain-text / markdown / docx parsers (unchanged)
# ---------------------------------------------------------------------------

def parse_txt(file_path: str) -> str:
    """Read a UTF-8 text file and return its contents as a string."""
    with open(file_path, "r", encoding="utf-8") as fh:
        return fh.read()


def parse_docx(file_path: str) -> str:
    """Read a .docx file with python-docx, joining non-empty paragraphs."""
    from docx import Document

    document = Document(file_path)
    paragraphs = [p.text for p in document.paragraphs if p.text]
    return "\n".join(paragraphs)


def parse_md(file_path: str) -> str:
    """Read a Markdown file as UTF-8 text (no structure parsing needed)."""
    with open(file_path, "r", encoding="utf-8") as fh:
        return fh.read()


# ---------------------------------------------------------------------------
# PDF helpers (V6-11)
# ---------------------------------------------------------------------------

def _detect_two_column(
    blocks: list[TextBlock], page_width: float
) -> tuple[bool, float]:
    """Detect whether a page uses a two-column layout.

    Returns ``(is_two_column, column_boundary)`` where *column_boundary*
    is the x-coordinate that separates left from right columns.
    """
    if not blocks or page_width <= 0:
        return False, page_width / 2

    mid_x = page_width / 2
    # A block belongs to the left column if it ends before mid_x;
    # to the right column if it starts after mid_x.
    left = [b for b in blocks if b.bbox[2] <= mid_x]
    right = [b for b in blocks if b.bbox[0] >= mid_x]

    if not left or not right:
        return False, mid_x

    # Among the narrow (non-spanning) blocks, find the largest gap in
    # center-x to determine the column boundary.
    centers = sorted(
        (b.bbox[0] + b.bbox[2]) / 2 for b in (left + right)
    )
    max_gap = 0.0
    boundary = mid_x
    for i in range(1, len(centers)):
        gap = centers[i] - centers[i - 1]
        if gap > max_gap:
            max_gap = gap
            boundary = (centers[i] + centers[i - 1]) / 2

    # The gap must be significant (> 15% of page width) to count as
    # two-column.
    if max_gap < page_width * 0.15:
        return False, mid_x

    return True, boundary


def _two_column_sort(
    blocks: list[TextBlock], page_width: float
) -> list[TextBlock]:
    """Sort blocks for two-column reading order.

    Left-column blocks (center_x < boundary) are read top-to-bottom first,
    then right-column blocks.  Full-width (spanning) blocks keep their
    vertical position and are emitted before/after the column region.
    """
    is_two_col, boundary = _detect_two_column(blocks, page_width)
    if not is_two_col:
        return sorted(blocks, key=lambda b: (b.bbox[1], b.bbox[0]))

    left = [b for b in blocks if (b.bbox[0] + b.bbox[2]) / 2 < boundary]
    right = [b for b in blocks if (b.bbox[0] + b.bbox[2]) / 2 >= boundary]

    if not left or not right:
        return sorted(blocks, key=lambda b: (b.bbox[1], b.bbox[0]))

    # Column region = intersection of the two columns' y-ranges.
    y_min = max(min(b.bbox[1] for b in left), min(b.bbox[1] for b in right))
    y_max = min(max(b.bbox[3] for b in left), max(b.bbox[3] for b in right))

    before = sorted(
        [b for b in blocks if b.bbox[1] < y_min],
        key=lambda b: (b.bbox[1], b.bbox[0]),
    )
    region_left = sorted(
        [b for b in left if y_min <= b.bbox[1] <= y_max],
        key=lambda b: (b.bbox[1], b.bbox[0]),
    )
    region_right = sorted(
        [b for b in right if y_min <= b.bbox[1] <= y_max],
        key=lambda b: (b.bbox[1], b.bbox[0]),
    )
    after = sorted(
        [b for b in blocks if b.bbox[1] > y_max],
        key=lambda b: (b.bbox[1], b.bbox[0]),
    )

    return before + region_left + region_right + after


def _font_stats_heading_threshold(blocks: list[TextBlock]) -> float:
    """Compute the heading font-size threshold from page statistics.

    Uses ``median + 1.0 * std`` so headings are detected relative to the
    body text on the page rather than by a hard-coded threshold.
    """
    sizes = [b.font_size for b in blocks if b.font_size and b.font_size > 0]
    if len(sizes) < 2:
        return 14.0  # fallback
    med = statistics.median(sizes)
    std = statistics.pstdev(sizes)
    return med + std


def _detect_footer_y_bands(
    all_page_blocks: list[list[TextBlock]], page_heights: list[float]
) -> set[int]:
    """Return the set of page indices that have a footer.

    A footer is a block in the bottom 10% of the page whose y-band
    (rounded) also appears on at least one other page.
    """
    # Collect y-bands at the bottom of each page.
    bottom_bands: dict[int, set[int]] = {}  # page_idx -> set of rounded y0
    for pi, blocks in enumerate(all_page_blocks):
        ph = page_heights[pi] if pi < len(page_heights) else 842.0
        bands = set()
        for b in blocks:
            if b.bbox[1] > ph * 0.90:  # bottom 10%
                bands.add(int(b.bbox[1] / 10) * 10)  # round to 10-unit band
        bottom_bands[pi] = bands

    # Cross-page repetition: a band appears on > 1 page.
    band_page_count: dict[int, int] = {}
    for bands in bottom_bands.values():
        for band in bands:
            band_page_count[band] = band_page_count.get(band, 0) + 1

    repeated_bands = {b for b, c in band_page_count.items() if c > 1}

    footer_pages: set[int] = set()
    for pi, bands in bottom_bands.items():
        if bands & repeated_bands:
            footer_pages.add(pi)
    return footer_pages


def _detect_table_grids(drawings: list) -> list[tuple]:
    """Detect table-grid bounding boxes from vector drawings.

    A table grid has >= 2 horizontal lines and >= 2 vertical lines that
    overlap in both x and y.  Returns a list of ``(x0, y0, x1, y1)``
    bounding boxes.
    """
    h_lines: list[tuple[float, float, float]] = []  # (x0, x1, y)
    v_lines: list[tuple[float, float, float]] = []  # (x, y0, y1)

    for d in drawings:
        rect = d.get("rect") or (0, 0, 0, 0)
        x0, y0, x1, y1 = rect
        width = x1 - x0
        height = y1 - y0
        if width > height * 5 and height < 3:
            # Horizontal line
            h_lines.append((x0, x1, (y0 + y1) / 2))
        elif height > width * 5 and width < 3:
            # Vertical line
            v_lines.append(((x0 + x1) / 2, y0, y1))

    if len(h_lines) < 2 or len(v_lines) < 2:
        return []

    # Find groups of horizontal lines that share a similar x-range.
    h_ys = sorted(set(int(h[2] / 5) * 5 for h in h_lines))
    v_xs = sorted(set(int(v[0] / 5) * 5 for v in v_lines))

    if len(h_ys) < 2 or len(v_xs) < 2:
        return []

    grid_x0 = min(v_xs)
    grid_x1 = max(v_xs)
    grid_y0 = min(h_ys)
    grid_y1 = max(h_ys)

    if grid_x1 > grid_x0 and grid_y1 > grid_y0:
        return [(grid_x0, grid_y0, grid_x1, grid_y1)]
    return []


def _detect_image_anchors(
    page: Any, page_no: int
) -> list[dict]:
    """Detect image/diagram anchors from raster images and filled shapes.

    Returns a list of dicts with ``bbox``, ``label``, ``is_decorative``.
    """
    anchors: list[dict] = []

    # Raster images.
    try:
        for img in page.get_images(full=True):
            xref = img[0]
            if xref == 0:
                continue
            try:
                rects = page.get_image_rects(xref)
            except Exception:
                rects = []
            for r in rects:
                area = (r.x1 - r.x0) * (r.y1 - r.y0)
                anchors.append({
                    "bbox": (r.x0, r.y0, r.x1, r.y1),
                    "label": "",
                    "is_decorative": area < 5000,
                })
    except Exception:
        pass

    # Filled-shape vector drawings (diagrams, icons).
    try:
        for d in page.get_drawings():
            rect = d.get("rect")
            if not rect:
                continue
            x0, y0, x1, y1 = rect
            if d.get("type") != "fs":
                continue
            # Only treat substantial filled shapes as image anchors
            # (skip tiny dots).
            area = (x1 - x0) * (y1 - y0)
            if area < 500:
                continue
            anchors.append({
                "bbox": (x0, y0, x1, y1),
                "label": "",
                "is_decorative": area < 10000,
            })
    except Exception:
        pass

    return anchors


# ---------------------------------------------------------------------------
# PDF parser (V6-11)
# ---------------------------------------------------------------------------

def parse_pdf(file_path: str) -> List[ParsedPage]:
    """Read a PDF file with PyMuPDF (fitz), returning ``[ParsedPage, ...]``.

    Page numbers are 1-indexed.  Each page's blocks are sorted using
    two-column-aware reading order, headings are detected from font-size
    statistics, footers from position + cross-page repetition, and tables
    from vector-drawing grids.
    """
    try:
        import fitz
    except ImportError:
        return _parse_pdf_pypdf_fallback(file_path)

    document = fitz.open(file_path)

    # --- First pass: extract raw blocks per page ---
    raw_pages: list[list[TextBlock]] = []
    page_heights: list[float] = []
    page_images: list[list[dict]] = []
    page_drawings: list[list] = []

    for page in document:
        ph = page.rect.height
        pw = page.rect.width
        page_heights.append(ph)

        blocks: list[TextBlock] = []
        for raw in page.get_text("dict").get("blocks", []):
            if raw.get("type") != 0:
                continue
            text = (
                "".join(
                    span.get("text", "")
                    for line in raw.get("lines", [])
                    for span in line.get("spans", [])
                ).strip()
            )
            if not text:
                continue
            x0, y0, x1, y1 = raw.get("bbox", (0, 0, 0, 0))
            size = max(
                (
                    float(span.get("size", 0))
                    for line in raw.get("lines", [])
                    for span in line.get("spans", [])
                ),
                default=0,
            )
            blocks.append(
                TextBlock(text, (x0, y0, x1, y1), size, "body", 0, "pdf")
            )

        raw_pages.append(blocks)
        page_images.append(_detect_image_anchors(page, len(raw_pages)))
        try:
            page_drawings.append(page.get_drawings())
        except Exception:
            page_drawings.append([])

    # --- Footer detection (needs all pages) ---
    footer_page_indices = _detect_footer_y_bands(raw_pages, page_heights)

    # --- Second pass: classify and order ---
    pages: List[ParsedPage] = []
    for pi, (raw_blocks, ph, imgs, drawings) in enumerate(
        zip(raw_pages, page_heights, page_images, page_drawings)
    ):
        pw = document[pi].rect.width

        # Table detection: blocks within drawing grids -> "table".
        table_grids = _detect_table_grids(drawings)
        classified = []
        for b in raw_blocks:
            bt = b.block_type
            for gx0, gy0, gx1, gy1 in table_grids:
                if gx0 - 5 <= b.bbox[0] and b.bbox[2] <= gx1 + 5 and gy0 - 5 <= b.bbox[1] and b.bbox[3] <= gy1 + 5:
                    bt = "table"
                    break
            if bt == "body" and pi in footer_page_indices and b.bbox[1] > ph * 0.90:
                bt = "footer"
            classified.append(
                TextBlock(b.text, b.bbox, b.font_size, bt, 0, "pdf", b.list_level)
            )

        # Heading detection from font statistics.
        heading_threshold = _font_stats_heading_threshold(classified)
        reclassified = []
        for b in classified:
            bt = b.block_type
            if bt == "body" and b.font_size and b.font_size > heading_threshold:
                bt = "title"
            reclassified.append(
                TextBlock(b.text, b.bbox, b.font_size, bt, 0, "pdf", b.list_level)
            )

        # Two-column-aware reading order.
        ordered = _two_column_sort(reclassified, pw)
        ordered = [
            TextBlock(
                b.text, b.bbox, b.font_size, b.block_type, order, b.source_kind, b.list_level
            )
            for order, b in enumerate(ordered)
        ]

        page_no = pi + 1
        # Page type: image_only if no text blocks but images exist.
        if not ordered and imgs:
            page_type = "image_only"
        elif not ordered:
            page_type = "image_only"
        else:
            page_type = "pdf"

        pages.append(
            ParsedPage(
                page_no=page_no,
                blocks=ordered,
                source_kind=page_type,
                parser_version="layout-v5",
                images=imgs,
                page_height=ph,
            )
        )

    return pages


def _parse_pdf_pypdf_fallback(file_path: str) -> List[ParsedPage]:
    """Fallback when PyMuPDF is not installed (uses pypdf)."""
    from pypdf import PdfReader

    return [
        ParsedPage(
            index + 1,
            [TextBlock(page.extract_text() or "", source_kind="pypdf")],
            "pdf",
            "pypdf-legacy",
        )
        for index, page in enumerate(PdfReader(file_path).pages)
    ]


# ---------------------------------------------------------------------------
# PPTX parser (V6-11)
# ---------------------------------------------------------------------------

def _is_title_placeholder(shape: Any) -> bool:
    """Return True if *shape* is a title or center-title placeholder."""
    if not getattr(shape, "is_placeholder", False):
        return False
    ph_type = getattr(shape.placeholder_format, "type", None)
    # 1 = TITLE, 3 = CENTER_TITLE (PP_PLACEHOLDER.TITLE / CENTER_TITLE)
    return ph_type in {1, 3}


def _is_image_only_slide(blocks: list[TextBlock], shapes: list) -> bool:
    """Detect a slide that is essentially an image with only a caption.

    Heuristic: the slide has shapes but no title heading, and all text
    blocks are short captions (start with 图/表/Figure or are < 20 chars).
    """
    if not blocks:
        return len(shapes) > 0
    has_heading = any(b.block_type == "title" for b in blocks)
    if has_heading:
        return False
    for b in blocks:
        stripped = b.text.strip()
        if len(stripped) > 20 and not stripped.startswith(("图", "表", "Figure", "Diagram")):
            return False
    return True


def parse_pptx(file_path: str) -> List[ParsedPage]:
    """Read a .pptx file with python-pptx, returning ``[ParsedPage, ...]``.

    Slide numbers are 1-indexed.  Shapes are sorted by ``(top, left)``
    for correct reading order.  Paragraph ``level`` is preserved as
    ``list_level``.  Title placeholders, body text, tables, and shapes
    are marked separately.  Single-letter shape labels are skipped.
    """
    from pptx import Presentation

    prs = Presentation(file_path)
    pages: List[ParsedPage] = []
    for index, slide in enumerate(prs.slides):
        blocks: list[TextBlock] = []
        shapes_list = list(slide.shapes)

        for shape in shapes_list:
            shape_left = float(shape.left) if shape.left is not None else 0
            shape_top = float(shape.top) if shape.top is not None else 0
            shape_width = float(shape.width) if shape.width is not None else 0
            shape_height = float(shape.height) if shape.height is not None else 0
            bbox = (
                shape_left,
                shape_top,
                shape_left + shape_width,
                shape_top + shape_height,
            )
            is_title = _is_title_placeholder(shape)

            # --- Text frame paragraphs ---
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = (para.text or "").strip()
                    if not text:
                        continue
                    # Skip single-letter shape labels (diagram annotations).
                    if len(text) <= 1:
                        continue
                    kind = "title" if is_title else "body"
                    level = getattr(para, "level", 0) or 0
                    blocks.append(
                        TextBlock(
                            text,
                            bbox,
                            None,
                            kind,
                            0,
                            "pptx",
                            level,
                        )
                    )

            # --- Table cells ---
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text = (cell.text or "").strip()
                        if not text:
                            continue
                        if len(text) <= 1:
                            continue
                        blocks.append(
                            TextBlock(
                                text,
                                bbox,
                                None,
                                "table",
                                0,
                                "pptx",
                                0,
                            )
                        )

        # Sort by (top, left) for reading order.
        ordered = sorted(blocks, key=lambda b: (b.bbox[1], b.bbox[0]))
        ordered = [
            TextBlock(
                b.text, b.bbox, b.font_size, b.block_type, order, b.source_kind, b.list_level
            )
            for order, b in enumerate(ordered)
        ]

        page_type = "image_only" if _is_image_only_slide(ordered, shapes_list) else "pptx"
        pages.append(ParsedPage(index + 1, ordered, page_type))

    return pages


# ---------------------------------------------------------------------------
# Dispatcher (unchanged)
# ---------------------------------------------------------------------------

def parse_file(
    file_path: str, file_type: str
) -> list:
    """Dispatch to the right parser by ``file_type`` (extension).

    Returns a list of ``(page_no, text)`` tuples. For non-paginated
    formats (txt, docx, md) ``page_no`` is ``None`` and the list has a
    single entry. Raises :class:`BusinessException` for unsupported types.
    """
    normalised = (file_type or "").lower().lstrip(".")
    if normalised == "txt":
        return [(None, parse_txt(file_path))]
    if normalised == "pdf":
        return parse_pdf(file_path)
    if normalised == "docx":
        return [(None, parse_docx(file_path))]
    if normalised == "md":
        return [(None, parse_md(file_path))]
    if normalised == "pptx":
        return parse_pptx(file_path)
    raise BusinessException(message=f"不支持的文件类型: {normalised}")
